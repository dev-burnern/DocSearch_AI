"""
Enhanced Text Extraction Module
다양한 문서 형식에서 텍스트 추출
- PDF (일반/스캔), Word, Excel, PPT, 이미지, HWP, 텍스트
- OCR 지원 (PaddleOCR)
- 테이블 추출 (pdfplumber)
"""
from __future__ import annotations

import io
import re
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import BinaryIO, Optional

from app.core.config import settings


@dataclass
class TextUnit:
    """Extracted text unit with metadata"""
    text: str
    page: int | None = None
    sheet: str | None = None
    slide: int | None = None
    heading: str | None = None
    is_table: bool = False
    is_ocr: bool = False
    metadata: dict = field(default_factory=dict)


@dataclass
class ExtractionResult:
    """Result of document extraction"""
    units: list[TextUnit]
    title: str | None = None
    author: str | None = None
    page_count: int | None = None
    metadata: dict = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)


# Regex patterns for text normalization
_CTRL_RE = re.compile(r"[\u0000-\u0008\u000B\u000C\u000E-\u001F]")
_ZW_RE = re.compile(r"[\u200B-\u200D\uFEFF]")
_MULTI_SPACE_RE = re.compile(r"[ \t]+")
_MULTI_NEWLINE_RE = re.compile(r"\n{3,}")


def normalize_text(s: str) -> str:
    """Normalize text by removing control characters and excessive whitespace"""
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    s = _CTRL_RE.sub(" ", s)
    s = _ZW_RE.sub("", s)
    s = _MULTI_SPACE_RE.sub(" ", s)
    s = _MULTI_NEWLINE_RE.sub("\n\n", s)
    return s.strip()


class TextExtractor:
    """Main text extractor class"""
    
    def __init__(self):
        self._ocr_model = None
    
    @property
    def ocr_model(self):
        """Lazy load OCR model"""
        if self._ocr_model is None and settings.ocr_enabled:
            try:
                from paddleocr import PaddleOCR
                self._ocr_model = PaddleOCR(
                    use_angle_cls=True,
                    lang='korean' if settings.ocr_lang == 'korean' else 'en',
                    use_gpu=settings.ocr_use_gpu,
                    show_log=False,
                )
            except ImportError:
                # Fallback to pytesseract
                pass
        return self._ocr_model
    
    def extract(
        self,
        file_path: Path | str,
        file_data: BinaryIO | bytes | None = None,
    ) -> ExtractionResult:
        """
        Extract text from a file
        
        Args:
            file_path: Path to the file or filename for type detection
            file_data: Optional file data (if not reading from disk)
        
        Returns:
            ExtractionResult with extracted text units
        """
        path = Path(file_path)
        ext = path.suffix.lower()
        
        # If file_data is provided, save to temp file for processing
        if file_data is not None:
            if isinstance(file_data, bytes):
                data = file_data
            else:
                data = file_data.read()
            
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(data)
                tmp_path = Path(tmp.name)
            
            try:
                return self._extract_by_type(tmp_path, ext)
            finally:
                tmp_path.unlink(missing_ok=True)
        else:
            return self._extract_by_type(path, ext)
    
    def _extract_by_type(self, path: Path, ext: str) -> ExtractionResult:
        """Route extraction based on file type"""
        extractors = {
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".xlsx": self._extract_xlsx,
            ".xlsm": self._extract_xlsx,
            ".pptx": self._extract_pptx,
            ".txt": self._extract_text,
            ".md": self._extract_text,
            ".hwp": self._extract_hwp,
            ".png": self._extract_image,
            ".jpg": self._extract_image,
            ".jpeg": self._extract_image,
            ".tiff": self._extract_image,
            ".bmp": self._extract_image,
        }
        
        extractor = extractors.get(ext)
        if not extractor:
            raise ValueError(f"Unsupported file type: {ext}")
        
        return extractor(path)
    
    def _extract_pdf(self, path: Path) -> ExtractionResult:
        """Extract text from PDF, with OCR fallback for scanned pages"""
        import fitz  # PyMuPDF
        
        units: list[TextUnit] = []
        warnings: list[str] = []
        
        doc = fitz.open(str(path))
        metadata = {
            "title": doc.metadata.get("title"),
            "author": doc.metadata.get("author"),
            "subject": doc.metadata.get("subject"),
            "creator": doc.metadata.get("creator"),
        }
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Try normal text extraction first
            text = page.get_text("text")
            
            if text.strip():
                units.append(TextUnit(
                    text=normalize_text(text),
                    page=page_num + 1,
                    is_ocr=False,
                ))
            elif settings.ocr_enabled:
                # Page might be scanned, try OCR
                ocr_text = self._ocr_page(page)
                if ocr_text:
                    units.append(TextUnit(
                        text=normalize_text(ocr_text),
                        page=page_num + 1,
                        is_ocr=True,
                    ))
                else:
                    warnings.append(f"Page {page_num + 1}: No text extracted (possibly blank or unreadable)")
        
        # Also extract tables
        table_units = self._extract_pdf_tables(path)
        units.extend(table_units)
        
        # Store page count before closing document
        page_count = len(doc)
        doc.close()
        
        if not units:
            raise RuntimeError(f"PDF text extraction produced empty result: {path}")
        
        return ExtractionResult(
            units=units,
            title=metadata.get("title"),
            author=metadata.get("author"),
            page_count=page_count,
            metadata=metadata,
            warnings=warnings,
        )
    
    def _extract_pdf_tables(self, path: Path) -> list[TextUnit]:
        """Extract tables from PDF using pdfplumber"""
        try:
            import pdfplumber
        except ImportError:
            return []
        
        units: list[TextUnit] = []
        
        with pdfplumber.open(str(path)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                
                for table_idx, table in enumerate(tables):
                    if not table:
                        continue
                    
                    # Convert table to markdown format
                    md_lines = []
                    headers = table[0] if table else []
                    
                    if headers:
                        header_row = "| " + " | ".join(str(h or "").strip() for h in headers) + " |"
                        md_lines.append(header_row)
                        md_lines.append("|" + "|".join(["---"] * len(headers)) + "|")
                    
                    for row in table[1:] if len(table) > 1 else []:
                        row_text = "| " + " | ".join(str(c or "").strip() for c in row) + " |"
                        md_lines.append(row_text)
                    
                    table_text = "\n".join(md_lines)
                    if table_text.strip():
                        units.append(TextUnit(
                            text=normalize_text(table_text),
                            page=page_num,
                            heading=f"Table {table_idx + 1}",
                            is_table=True,
                        ))
        
        return units
    
    def _ocr_page(self, page) -> str:
        """OCR a PDF page"""
        if not self.ocr_model:
            return ""
        
        # Render page to image
        pix = page.get_pixmap(dpi=300)
        img_data = pix.tobytes("png")
        
        # Save to temp file for OCR
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(img_data)
            tmp_path = tmp.name
        
        try:
            result = self.ocr_model.ocr(tmp_path, cls=True)
            if result and result[0]:
                texts = [line[1][0] for line in result[0] if line[1]]
                return "\n".join(texts)
        finally:
            Path(tmp_path).unlink(missing_ok=True)
        
        return ""
    
    def _extract_docx(self, path: Path) -> ExtractionResult:
        """Extract text from Word document"""
        from docx import Document
        from docx.opc.exceptions import PackageNotFoundError
        
        try:
            doc = Document(str(path))
        except PackageNotFoundError:
            raise RuntimeError(f"Invalid or corrupted DOCX file: {path}")
        
        units: list[TextUnit] = []
        current_heading: str | None = None
        current_text: list[str] = []
        
        def flush_section():
            nonlocal current_text, current_heading
            if current_text:
                text = "\n\n".join(current_text)
                if text.strip():
                    units.append(TextUnit(
                        text=normalize_text(text),
                        heading=current_heading,
                    ))
                current_text = []
        
        # Extract paragraphs with heading tracking
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                flush_section()
                current_heading = para.text.strip()
            else:
                text = para.text.strip()
                if text:
                    current_text.append(text)
        
        flush_section()
        
        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            
            if rows:
                # Create markdown table
                if len(rows) > 1:
                    md_lines = [
                        "| " + rows[0] + " |",
                        "|" + "|".join(["---"] * rows[0].count("|")) + "|",
                    ]
                    md_lines.extend(["| " + r + " |" for r in rows[1:]])
                    table_text = "\n".join(md_lines)
                else:
                    table_text = rows[0]
                
                units.append(TextUnit(
                    text=normalize_text(table_text),
                    heading=f"Table {table_idx + 1}",
                    is_table=True,
                ))
        
        # Get metadata
        core_props = doc.core_properties
        metadata = {
            "title": core_props.title,
            "author": core_props.author,
            "subject": core_props.subject,
            "created": str(core_props.created) if core_props.created else None,
        }
        
        if not units:
            raise RuntimeError(f"DOCX text extraction produced empty result: {path}")
        
        return ExtractionResult(
            units=units,
            title=metadata.get("title"),
            author=metadata.get("author"),
            metadata=metadata,
        )
    
    def _extract_xlsx(self, path: Path) -> ExtractionResult:
        """Extract text from Excel workbook"""
        import openpyxl
        
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        units: list[TextUnit] = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            
            for row in ws.iter_rows(values_only=True):
                vals = []
                for v in row:
                    if v is not None:
                        s = str(v).strip()
                        if s:
                            vals.append(s)
                if vals:
                    rows.append(" | ".join(vals))
            
            if rows:
                # Create markdown table format
                text = "\n".join(rows)
                units.append(TextUnit(
                    text=normalize_text(text),
                    sheet=sheet_name,
                    is_table=True,
                ))
        
        wb.close()
        
        if not units:
            raise RuntimeError(f"XLSX text extraction produced empty result: {path}")
        
        return ExtractionResult(units=units)
    
    def _extract_pptx(self, path: Path) -> ExtractionResult:
        """Extract text from PowerPoint presentation"""
        from pptx import Presentation
        
        prs = Presentation(str(path))
        units: list[TextUnit] = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                
                # Extract table content
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells]
                        if any(row_texts):
                            texts.append(" | ".join(row_texts))
            
            # Extract notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"\n[노트]\n{notes}")
            
            if texts:
                units.append(TextUnit(
                    text=normalize_text("\n\n".join(texts)),
                    slide=slide_num,
                ))
        
        # Get metadata
        core_props = prs.core_properties
        metadata = {
            "title": core_props.title,
            "author": core_props.author,
            "subject": core_props.subject,
        }
        
        if not units:
            raise RuntimeError(f"PPTX text extraction produced empty result: {path}")
        
        return ExtractionResult(
            units=units,
            title=metadata.get("title"),
            author=metadata.get("author"),
            page_count=len(prs.slides),
            metadata=metadata,
        )
    
    def _extract_text(self, path: Path) -> ExtractionResult:
        """Extract text from plain text file"""
        text = path.read_text(encoding="utf-8", errors="ignore")
        text = normalize_text(text)
        
        if not text:
            raise RuntimeError(f"Text file is empty: {path}")
        
        return ExtractionResult(
            units=[TextUnit(text=text)],
        )
    
    def _extract_hwp(self, path: Path) -> ExtractionResult:
        """Extract text from HWP (Hangul Word Processor) file"""
        try:
            import olefile
            import zlib
        except ImportError:
            raise RuntimeError("olefile package required for HWP extraction")
        
        texts: list[str] = []
        
        try:
            ole = olefile.OleFileIO(str(path))
            
            # Try to extract from BodyText sections
            for stream in ole.listdir():
                if stream[0] == "BodyText":
                    data = ole.openstream(stream).read()
                    
                    # Try to decompress
                    try:
                        decompressed = zlib.decompress(data, -15)
                        # HWP uses UTF-16 LE encoding
                        text = decompressed.decode("utf-16-le", errors="ignore")
                        # Remove control characters
                        text = normalize_text(text)
                        if text:
                            texts.append(text)
                    except zlib.error:
                        # Not compressed, try direct decode
                        text = data.decode("utf-16-le", errors="ignore")
                        text = normalize_text(text)
                        if text:
                            texts.append(text)
            
            ole.close()
        except Exception as e:
            raise RuntimeError(f"Failed to extract HWP: {e}")
        
        if not texts:
            raise RuntimeError(f"HWP text extraction produced empty result: {path}")
        
        return ExtractionResult(
            units=[TextUnit(text="\n\n".join(texts))],
            warnings=["HWP extraction is limited; some formatting may be lost"],
        )
    
    def _extract_image(self, path: Path) -> ExtractionResult:
        """Extract text from image using OCR"""
        if not self.ocr_model:
            raise RuntimeError("OCR is not enabled or available")
        
        result = self.ocr_model.ocr(str(path), cls=True)
        
        if result and result[0]:
            texts = [line[1][0] for line in result[0] if line[1]]
            text = "\n".join(texts)
            
            if text.strip():
                return ExtractionResult(
                    units=[TextUnit(text=normalize_text(text), is_ocr=True)],
                )
        
        raise RuntimeError(f"Image OCR produced empty result: {path}")


# Singleton instance
_extractor: TextExtractor | None = None


def get_text_extractor() -> TextExtractor:
    """Get text extractor singleton"""
    global _extractor
    if _extractor is None:
        _extractor = TextExtractor()
    return _extractor


def extract_text_from_file(
    file_path: Path | str,
    file_data: BinaryIO | bytes | None = None,
) -> ExtractionResult:
    """Convenience function to extract text from a file"""
    return get_text_extractor().extract(file_path, file_data)
